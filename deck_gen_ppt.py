from pptx import Presentation  # python-pptx
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class PPTGenerator:
    def __init__(self, plots, templates):
        self.plots_folder = plots
        self.ppt = Presentation(templates)


    def page1_group(self, group):
        textbox = self.ppt.slides[0].shapes.add_textbox(Inches(0.68), Inches(1.88), height=Inches(0.49), width=Inches(8.01))
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = group
        font = run.font
        font.name = 'Arial'
        font.size = Pt(28)
        font.bold = False
        font.italic = False
        font.color.rgb = RGBColor(0, 41, 96)


    def page1_title(self, title):
        textbox = self.ppt.slides[0].shapes.add_textbox(Inches(0.7), Inches(2.39), height=Inches(0.42), width=Inches(8.01))
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = title
        font = run.font
        font.name = 'Arial'
        font.size = Pt(24)
        font.bold = False
        font.italic = False
        font.color.rgb = RGBColor(0, 96, 169)


    def page1_date(self, month, date, year):
        textbox = self.ppt.slides[0].shapes.add_textbox(Inches(0.82), Inches(6.46), height=Inches(0.26), width=Inches(3.34))
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT

        run = paragraph.add_run()
        run.text = f'{month}{date[0]}'
        font = run.font
        font.name = 'Arial'
        font.size = Pt(14)
        font.bold = True
        font.italic = False
        font.color.rgb = RGBColor(0, 0, 0)

        suffix_run = paragraph.add_run()
        suffix_run.text = date[1]  # suffix
        suffix_font = suffix_run.font
        suffix_font.name = font.name
        suffix_font.size = font.size
        suffix_font.bold = font.bold
        suffix_font.italic = font.italic
        suffix_font.color.rgb = font.color.rgb

        year_run = paragraph.add_run()
        year_run.text = f', {year}'
        year_font = year_run.font
        year_font.name = font.name
        year_font.size = font.size
        year_font.bold = font.bold
        year_font.italic = font.italic
        year_font.color.rgb = font.color.rgb


    def page2_title(self, title):
        textbox = self.ppt.slides[1].shapes.add_textbox(Inches(0.58), Inches(0.55), height=Inches(0.63), width=Inches(11.08))
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = title
        font = run.font
        font.name = 'Arial'
        font.size = Pt(24)
        font.bold = False
        font.italic = False
        font.color.rgb = RGBColor(0, 96, 169)


    def page2_plot_title(self, title):
        textbox = self.ppt.slides[1].shapes.add_textbox(Inches(1.28), Inches(1.55), height=Inches(0.26), width=Inches(6.01))
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = title
        font = run.font
        font.name = 'Arial'
        font.size = Pt(14)
        font.bold = True
        font.italic = False
        font.color.rgb = RGBColor(255, 255, 255)


    def page2_plot(self, group):
        self.ppt.slides[1].shapes.add_picture(f'{self.plots_folder}{group}_comp.png', Inches(0.72), Inches(2), height=Inches(3.35), width=Inches(7.11))


    def page2_table(self, table):
        self.ppt.slides[1].shapes.add_picture(table, Inches(0.46), Inches(5.69), height=Inches(0.98), width=Inches(7.63))


    def page2_commentary_1(self, group, comments):
        if group == 'Foreign':
            self.ppt.slides[1].shapes.add_picture(comments, Inches(8.8), Inches(2.3), height=Inches(1.04), width=Inches(4.14))
        else:
            self.ppt.slides[1].shapes.add_picture(comments, Inches(8.8), Inches(2.3), height=Inches(1.13), width=Inches(4.1))


    def page2_commentary_2(self, comments):  # Foreign only
        self.ppt.slides[1].shapes.add_picture(comments, Inches(8.8), Inches(4.4), height=Inches(0.58), width=Inches(4.14))


    def page3_title(self, title):
        textbox = self.ppt.slides[2].shapes.add_textbox(Inches(0.58), Inches(0.55), height=Inches(0.63), width=Inches(11.08))
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = title
        font = run.font
        font.name = 'Arial'
        font.size = Pt(24)
        font.bold = False
        font.italic = False
        font.color.rgb = RGBColor(0, 96, 169)

    
    def page3_table(self, table):
        self.ppt.slides[2].shapes.add_picture(table, Inches(1.3), Inches(2.5), height=Inches(4.35), width=Inches(11.04))


    def page4_title(self, title):  # Local only
        textbox = self.ppt.slides[3].shapes.add_textbox(Inches(0.58), Inches(0.55), height=Inches(0.63), width=Inches(11.08))
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = title
        font = run.font
        font.name = 'Arial'
        font.size = Pt(24)
        font.bold = False
        font.italic = False
        font.color.rgb = RGBColor(0, 96, 169)


    def page4_table(self, table):  # Local only
        self.ppt.slides[3].shapes.add_picture(table, Inches(0.85), Inches(1.41), height=Inches(1.09), width=Inches(11.63))


    def page4_plot_1(self, group):  # Local only
        self.ppt.slides[3].shapes.add_picture(f'{self.plots_folder}{group}_shrk_pie.png', Inches(0.37), Inches(2.55), height=Inches(3.3), width=Inches(6.1))


    def page4_plot_2(self, group):
        self.ppt.slides[3].shapes.add_picture(f'{self.plots_folder}{group}_shrk_bar.png', Inches(6.87), Inches(2.55), height=Inches(3.3), width=Inches(6.1))


    def page4_delete(self):  # select groups only
        self.ppt.slides._sldIdLst.remove(self.ppt.slides._sldIdLst[3])
    

    def save_ppt(self, path, year, month, title):
        self.ppt.save(f'{path}\\{title}_Workforce_Report_{month}{year}.pptx')
